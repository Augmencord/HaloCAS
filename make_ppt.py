from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def main():
    prs = Presentation()
    
    # Common layout references
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    title_only_layout = prs.slide_layouts[5]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "HaloCAS"
    subtitle.text = "Next-Generation Safety for Heavy Machinery\nPowered by Intelligent Graph Workflows"

    # Slide 2: The Problem (with Image 1)
    slide = prs.slides.add_slide(title_only_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "The Problem: Blind Spots in the Danger Zone"
    
    img_path1 = r"C:\Users\augme\.gemini\antigravity-ide\brain\0e4b0bcf-a091-4dda-accc-efeb7cdfceeb\industrial_mining_safety_1786403179215.png"
    if os.path.exists(img_path1):
        shapes.add_picture(img_path1, Inches(0.5), Inches(1.5), width=Inches(4.5))
        
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "In mining and manufacturing, humans and colossal machines operate side-by-side."
    
    p2 = tf.add_paragraph()
    p2.text = "• The Risk: Heavy machinery creates deadly blind spots."
    p2.level = 0
    p3 = tf.add_paragraph()
    p3.text = "• The Flaw: Standard cameras require constant monitoring."
    p3.level = 0
    p4 = tf.add_paragraph()
    p4.text = "• The Result: 'Alarm fatigue' from sensors that beep at everything."
    p4.level = 0

    # Slide 3: The Solution / Graph Workflow (with Image 2)
    slide = prs.slides.add_slide(title_only_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "The Solution: Intelligent Graph Workflow"
    
    img_path2 = r"C:\Users\augme\.gemini\antigravity-ide\brain\0e4b0bcf-a091-4dda-accc-efeb7cdfceeb\halocas_graph_workflow_1786403194640.png"
    if os.path.exists(img_path2):
        shapes.add_picture(img_path2, Inches(5.0), Inches(1.5), width=Inches(4.5))

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.0), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HaloCAS doesn't just beep—it thinks."
    
    p2 = tf.add_paragraph()
    p2.text = "1. State Check: Is the machine moving?"
    p3 = tf.add_paragraph()
    p3.text = "2. Proximity: Has a worker breached 10m?"
    p4 = tf.add_paragraph()
    p4.text = "3. Trajectory: Is the machine moving towards the worker?"
    p5 = tf.add_paragraph()
    p5.text = "4. Auth Scan: Is it an authorized mechanic?"

    # Slide 4: Actionable UX (with Image 3)
    slide = prs.slides.add_slide(title_only_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Actionable User Experience"
    
    img_path3 = r"C:\Users\augme\.gemini\antigravity-ide\brain\0e4b0bcf-a091-4dda-accc-efeb7cdfceeb\halocas_alert_dashboard_1786403206734.png"
    if os.path.exists(img_path3):
        shapes.add_picture(img_path3, Inches(0.5), Inches(1.5), width=Inches(5.0))
        
    txBox = slide.shapes.add_textbox(Inches(5.8), Inches(1.5), Inches(3.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HaloCAS operates completely invisibly to operators on the ground."
    
    p2 = tf.add_paragraph()
    p2.text = "When an incident occurs, Safety Officers receive:"
    p3 = tf.add_paragraph()
    p3.text = "• Instant automated email alerts."
    p4 = tf.add_paragraph()
    p4.text = "• Clear context of who and when."
    p5 = tf.add_paragraph()
    p5.text = "• Attached MP4 video clips for undeniable evidence."

    # Slide 5: Code / GitHub
    slide = prs.slides.add_slide(title_only_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "See the Code"
    
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Check out the open-source code and technical architecture on GitHub:\n\nhttps://github.com/Augmencord/HaloCAS"
    
    prs.save(r"C:\Users\augme\.gemini\antigravity-ide\scratch\HaloCAS\HaloCAS_Presentation_v2.pptx")
    print("Presentation generated successfully!")

if __name__ == '__main__':
    main()
